"""
generate_pptx.py — Professional dark + teal presentation
Style: Dark background, teal/cyan accents, clean cards, modern feel
Run: python -m pip install python-pptx
Then: python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# =============================================================================
# COLOR SCHEME — Dark + Teal (matching the reference template)
# =============================================================================
BG          = RGBColor(0x11, 0x14, 0x1A)    # near-black background
BG_CARD     = RGBColor(0x1A, 0x1F, 0x2B)    # card background
BG_CARD2    = RGBColor(0x21, 0x27, 0x35)    # lighter card
BG_INNER    = RGBColor(0x15, 0x19, 0x23)    # inner/nested card
TEAL        = RGBColor(0x2D, 0xD4, 0xBF)    # primary teal accent
TEAL_DIM    = RGBColor(0x1A, 0x7A, 0x6E)    # darker teal
TEAL_DARK   = RGBColor(0x0F, 0x3D, 0x38)    # very dark teal for borders
CYAN        = RGBColor(0x67, 0xE8, 0xF9)    # bright cyan highlight
WHITE       = RGBColor(0xF1, 0xF5, 0xF9)    # near-white text
GRAY        = RGBColor(0x94, 0xA3, 0xB8)    # secondary text
DIM         = RGBColor(0x64, 0x74, 0x8B)    # dimmed text
RED         = RGBColor(0xF8, 0x71, 0x71)    # red/warning
YELLOW      = RGBColor(0xFA, 0xCC, 0x15)    # yellow highlight
BLUE        = RGBColor(0x60, 0xA5, 0xFA)    # blue
PURPLE      = RGBColor(0xA7, 0x8B, 0xFA)    # purple
ORANGE      = RGBColor(0xFB, 0x92, 0x3C)    # orange

TOTAL_SLIDES = 9  # added agenda slide


def set_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG


def bar(slide):
    """Top accent bar."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.05))
    s.fill.solid(); s.fill.fore_color.rgb = TEAL; s.line.fill.background()


def circle_deco(slide, x, y, size, color=TEAL_DARK):
    """Decorative background circle."""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def card(slide, l, t, w, h, color=BG_CARD, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.adjustments[0] = 0.06
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text="", size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb, tf


def para(tf, text="", size=12, color=WHITE, bold=False, space=0, align=PP_ALIGN.LEFT, font="Calibri"):
    p = tf.add_paragraph(); p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold; p.font.name = font
    p.space_before = Pt(space); p.alignment = align
    return p


def footer(slide, n):
    txt(slide, 0.4, 7.1, 4, 0.3, "Predictive Code-Switching Detection", 8, DIM)
    txt(slide, 8.5, 7.1, 1.2, 0.3, f"{n} / {TOTAL_SLIDES}", 8, DIM, align=PP_ALIGN.RIGHT)


def teal_bullet(slide, x, y, text_str, size=10):
    """Teal circle bullet + text."""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.07), Inches(0.1), Inches(0.1))
    s.fill.solid(); s.fill.fore_color.rgb = TEAL; s.line.fill.background()
    txt(slide, x + 0.2, y, 3.8, 0.25, text_str, size, GRAY)


# =============================================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# =============================================================================
# SLIDE 1: TITLE
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)
circle_deco(sl, -0.5, 4.5, 3)
circle_deco(sl, 7.5, 0.5, 4, RGBColor(0x0A, 0x2A, 0x28))
circle_deco(sl, 1, 0.8, 1.5, RGBColor(0x0D, 0x33, 0x30))

txt(sl, 1.5, 2.0, 7, 0.8, "Predictive Code-Switching", 38, WHITE, True, PP_ALIGN.CENTER)
txt(sl, 1.5, 2.8, 7, 0.6, "Detection", 38, TEAL, True, PP_ALIGN.CENTER)

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.6), Inches(2), Inches(0.04))
s.fill.solid(); s.fill.fore_color.rgb = TEAL; s.line.fill.background()

txt(sl, 1.5, 3.9, 7, 0.4, "A Universality-Focused Multitask Learning Approach", 14, GRAY, align=PP_ALIGN.CENTER)
txt(sl, 1.5, 4.6, 7, 0.4, "Project Update 1: Streaming Data & Causal Baselines", 12, TEAL_DIM, align=PP_ALIGN.CENTER)

# Team info cards
for i, name in enumerate(["Jae Hun Cho", "Fengying Zeng", "Xiaoyan Cai"]):
    x = 2.2 + i * 2.0
    card(sl, x, 5.3, 1.8, 0.6, BG_CARD, TEAL_DARK)
    txt(sl, x, 5.4, 1.8, 0.4, name, 10, TEAL, align=PP_ALIGN.CENTER)

txt(sl, 1.5, 6.1, 7, 0.3, "Northeastern University, Seattle  ·  February 25, 2026", 10, DIM, align=PP_ALIGN.CENTER)
footer(sl, 1)


# =============================================================================
# SLIDE 2: AGENDA
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)
circle_deco(sl, 7, 5, 3.5, RGBColor(0x0A, 0x2A, 0x28))

txt(sl, 0.5, 0.4, 4, 0.6, "Agenda", 30, WHITE, True)
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.2), Inches(0.04))
s.fill.solid(); s.fill.fore_color.rgb = TEAL; s.line.fill.background()

agenda_items = [
    ("01", "Language Pair Selection", "6 pairs across 2 diversity axes", "Slide 3"),
    ("02", "Data Exploration", "49,833 samples — statistics & distributions", "Slide 4"),
    ("03", "Streaming DataLoader", "On-the-fly pipeline from HuggingFace to batches", "Slide 5"),
    ("04", "Shifted Label Generation", "y_sw (switch) and y_dur (duration) creation", "Slide 6"),
    ("05", "Causal Attention Mask", "Prefix-only context enforcement", "Slide 7"),
    ("06", "Baseline Metrics", "Naive baseline F1 & universality σ", "Slide 8"),
    ("07", "Challenges & Next Steps", "Sub-word issues and Phase 2 roadmap", "Slide 9"),
]

for i, (num, title, desc, page) in enumerate(agenda_items):
    y = 1.4 + i * 0.78
    card(sl, 0.5, y, 8.5, 0.65, BG_CARD, TEAL_DARK if i % 2 == 0 else None)
    
    # Number circle
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y + 0.1), Inches(0.45), Inches(0.45))
    s.fill.solid(); s.fill.fore_color.rgb = TEAL_DARK; s.line.color.rgb = TEAL; s.line.width = Pt(1)
    txt(sl, 0.7, y + 0.12, 0.45, 0.4, num, 13, TEAL, True, PP_ALIGN.CENTER)
    
    txt(sl, 1.35, y + 0.08, 4.5, 0.25, title, 13, WHITE, True)
    txt(sl, 1.35, y + 0.35, 4.5, 0.25, desc, 9, DIM)
    txt(sl, 7.8, y + 0.18, 1, 0.25, page, 9, TEAL_DIM, align=PP_ALIGN.RIGHT)

footer(sl, 2)


# =============================================================================
# SLIDE 3: LANGUAGE PAIR SELECTION
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)
circle_deco(sl, -1, 5.5, 3, RGBColor(0x0A, 0x2A, 0x28))

txt(sl, 0.5, 0.4, 9, 0.6, "Language Pair Selection", 28, WHITE, True)
txt(sl, 0.5, 0.95, 9, 0.35, "6 pairs across 2 diversity axes to rigorously test universality", 12, DIM)

# Axis 1
card(sl, 0.4, 1.5, 4.4, 3.3, BG_CARD, TEAL_DIM)
s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(1.6), Inches(0.25), Inches(0.25))
s.fill.solid(); s.fill.fore_color.rgb = TEAL; s.line.fill.background()
txt(sl, 0.95, 1.6, 3.5, 0.3, "Axis 1: Script Difference", 14, TEAL, True)
txt(sl, 0.6, 2.0, 4, 0.25, "Can the model work without character-level cues?", 9, DIM)

card(sl, 0.6, 2.4, 4, 0.95, BG_INNER)
txt(sl, 0.75, 2.45, 3.7, 0.25, "✦ Different Script — Easy LID", 10, TEAL, True)
txt(sl, 0.75, 2.75, 3.7, 0.2, "Hindi-English  (Devanagari ↔ Latin)", 9, GRAY)
txt(sl, 0.75, 2.97, 3.7, 0.2, "Arabic-English  (Arabic ↔ Latin)", 9, GRAY)

card(sl, 0.6, 3.5, 4, 0.95, BG_INNER)
txt(sl, 0.75, 3.55, 3.7, 0.25, "✦ Same Script — Hard LID", 10, RED, True)
txt(sl, 0.75, 3.85, 3.7, 0.2, "Spanish-English  (Latin ↔ Latin)", 9, GRAY)
txt(sl, 0.75, 4.07, 3.7, 0.2, "French-English  (Latin ↔ Latin)", 9, GRAY)

# Axis 2
card(sl, 5.2, 1.5, 4.4, 3.3, BG_CARD, PURPLE)
s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.4), Inches(1.6), Inches(0.25), Inches(0.25))
s.fill.solid(); s.fill.fore_color.rgb = PURPLE; s.line.fill.background()
txt(sl, 5.75, 1.6, 3.5, 0.3, "Axis 2: Typological Distance", 14, PURPLE, True)
txt(sl, 5.4, 2.0, 4, 0.25, "Can the model generalize across different grammars?", 9, DIM)

card(sl, 5.4, 2.4, 4, 0.95, BG_INNER)
txt(sl, 5.55, 2.45, 3.7, 0.25, "✦ Close (SVO) — Similar word order", 10, TEAL, True)
txt(sl, 5.55, 2.75, 3.7, 0.2, "Spanish-EN, French-EN", 9, GRAY)
txt(sl, 5.55, 2.97, 3.7, 0.2, "Both Subject-Verb-Object", 9, GRAY)

card(sl, 5.4, 3.5, 4, 0.95, BG_INNER)
txt(sl, 5.55, 3.55, 3.7, 0.25, "✦ Distant (SOV / Topic-prominent)", 10, RED, True)
txt(sl, 5.55, 3.85, 3.7, 0.2, "Korean-English  (SOV ↔ SVO)", 9, GRAY)
txt(sl, 5.55, 4.07, 3.7, 0.2, "Chinese-English  (Topic-prominent)", 9, GRAY)

# Insight box
card(sl, 0.4, 5.1, 9.2, 0.7, BG_CARD, TEAL_DARK)
_, tf = txt(sl, 0.6, 5.2, 8.8, 0.5, "", 10)
r = tf.paragraphs[0].add_run()
r.text = "💡 Key insight: "; r.font.color.rgb = YELLOW; r.font.bold = True; r.font.size = Pt(10)
r2 = tf.paragraphs[0].add_run()
r2.text = "All 6 pairs have ~8K samples with similar quality scores (~8.5/10), eliminating data quantity as a confound. Any performance gap = genuine linguistic difficulty."
r2.font.color.rgb = GRAY; r2.font.size = Pt(9)

footer(sl, 3)


# =============================================================================
# SLIDE 4: DATA EXPLORATION
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)

txt(sl, 0.5, 0.4, 9, 0.6, "Data Exploration Results", 28, WHITE, True)
txt(sl, 0.5, 0.95, 9, 0.35, "SwitchLingua — 49,833 samples across our 6 selected pairs", 12, DIM)

# Table
tbl_shape = sl.shapes.add_table(8, 6, Inches(0.4), Inches(1.4), Inches(9.2), Inches(2.8))
tbl = tbl_shape.table

headers = ["Language Pair", "Samples", "Share", "Avg Score", "Avg Words", "Scripts"]
rows = [
    ["Hindi-English", "8,381", "16.8%", "8.49", "56.8", "Devanagari ↔ Latin"],
    ["Arabic-English", "8,548", "17.2%", "8.44", "45.9", "Arabic ↔ Latin"],
    ["Spanish-English", "8,110", "16.3%", "8.54", "54.9", "Latin ↔ Latin"],
    ["French-English", "8,070", "16.2%", "8.47", "53.3", "Latin ↔ Latin"],
    ["Korean-English", "8,266", "16.6%", "8.55", "44.8", "Hangul ↔ Latin"],
    ["Chinese-English", "8,458", "17.0%", "8.49", "27.5", "CJK ↔ Latin"],
    ["TOTAL", "49,833", "100%", "~8.5", "—", ""],
]
pair_colors = [CYAN, CYAN, ORANGE, ORANGE, PURPLE, PURPLE, WHITE]

for j, h in enumerate(headers):
    c = tbl.cell(0, j); c.text = h
    p = c.text_frame.paragraphs[0]; p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
    c.fill.solid(); c.fill.fore_color.rgb = TEAL_DARK

for i, row in enumerate(rows):
    for j, val in enumerate(row):
        c = tbl.cell(i+1, j); c.text = val
        p = c.text_frame.paragraphs[0]; p.font.size = Pt(9); p.font.name = "Calibri"
        p.font.color.rgb = pair_colors[i] if j == 0 else GRAY
        p.font.bold = (j == 0 or i == 6)
        c.fill.solid(); c.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else BG_INNER
        if i == 6: c.fill.fore_color.rgb = TEAL_DARK

# Stats boxes
stats = [("1.1x", "Balance Ratio", "Well balanced ✓", TEAL),
         ("~⅓ each", "CS Type Dist.", "Intra / Inter / Tag", CYAN),
         ("19.1%", "Switch Rate", "Class imbalance 1:4", YELLOW),
         ("27.5w", "Chinese-EN Length", "Notably shorter ⚠", ORANGE)]
for i, (val, lbl, sub, col) in enumerate(stats):
    x = 0.4 + i * 2.35
    card(sl, x, 4.5, 2.15, 1.0, BG_CARD, TEAL_DARK)
    txt(sl, x, 4.55, 2.15, 0.35, val, 20, col, True, PP_ALIGN.CENTER)
    txt(sl, x, 4.95, 2.15, 0.2, lbl, 8, DIM, align=PP_ALIGN.CENTER)
    txt(sl, x, 5.15, 2.15, 0.2, sub, 8, col, align=PP_ALIGN.CENTER)

footer(sl, 4)


# =============================================================================
# SLIDE 5: STREAMING DATALOADER
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)

txt(sl, 0.5, 0.4, 9, 0.6, "Streaming DataLoader", 28, WHITE, True)
txt(sl, 0.5, 0.95, 9, 0.35, "On-the-fly processing — never loads full 420K samples into memory", 12, DIM)

# Pipeline
card(sl, 0.4, 1.4, 9.2, 1.2, BG_CARD, TEAL_DARK)
steps = [("📦", "HuggingFace\nCSV Files"), ("🔍", "Filter Pair\n& Score"), ("🏷️", "Token-Level\nLanguage ID"),
         ("🔄", "Shifted\nLabels"), ("✂️", "Subword\nAlignment"), ("📊", "Batch\nYield")]
for i, (icon, label) in enumerate(steps):
    x = 0.65 + i * 1.52
    txt(sl, x, 1.5, 0.4, 0.35, icon, 18, WHITE, align=PP_ALIGN.CENTER)
    txt(sl, x - 0.3, 1.85, 1.0, 0.6, label, 8, GRAY, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(sl, x + 0.6, 1.6, 0.5, 0.3, "→", 16, TEAL_DIM, align=PP_ALIGN.CENTER)

# Left card
card(sl, 0.4, 2.9, 4.4, 2.6, BG_CARD)
txt(sl, 0.6, 3.0, 4, 0.3, "Why Streaming?", 14, TEAL, True)
bullets = ["Memory efficient — processes one sample at a time",
           "Per-file loading via data_files= parameter",
           "Round-robin interleaving for mixed batches",
           "Simulates real-time keyboard input scenario"]
for i, b in enumerate(bullets):
    teal_bullet(sl, 0.7, 3.45 + i * 0.38, b, 9)

# Right card — code
card(sl, 5.2, 2.9, 4.4, 2.6, BG_CARD)
txt(sl, 5.4, 3.0, 4, 0.3, "Key Implementation", 14, TEAL, True)
card(sl, 5.4, 3.4, 4.0, 1.9, BG_INNER)
code = ('# Load specific pair file\n'
        'ds = load_dataset(\n'
        '    "SwitchLingua_text",\n'
        '    data_files="Hindi_eng.csv",\n'
        '    streaming=True\n'
        ')\n\n'
        '# PyTorch IterableDataset\n'
        'class StreamingCSDataset(\n'
        '    IterableDataset):\n'
        '  def __iter__(self):\n'
        '    yield sample')
txt(sl, 5.5, 3.45, 3.8, 1.8, code, 7, GRAY, font="Consolas")

# Output bar
card(sl, 0.4, 5.7, 9.2, 0.45, BG_CARD, TEAL_DARK)
_, tf = txt(sl, 0.6, 5.75, 8.8, 0.35, "", 9)
r = tf.paragraphs[0].add_run(); r.text = "Output: "; r.font.color.rgb = TEAL; r.font.bold = True; r.font.size = Pt(9)
r2 = tf.paragraphs[0].add_run()
r2.text = "input_ids (B,256) · attention_mask (B,256) · switch_labels (B,256) · duration_labels (B,256)"
r2.font.color.rgb = CYAN; r2.font.size = Pt(8); r2.font.name = "Consolas"

footer(sl, 5)


# =============================================================================
# SLIDE 6: SHIFTED LABELS
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)

txt(sl, 0.5, 0.4, 9, 0.6, "Shifted Label Generation", 28, WHITE, True)
txt(sl, 0.5, 0.95, 9, 0.35, "At position t, predict what happens at position t+1", 12, DIM)

# Example table
card(sl, 0.4, 1.35, 9.2, 2.4, BG_CARD, TEAL_DARK)
txt(sl, 0.6, 1.4, 4, 0.25, "Example: Hindi-English sentence", 11, YELLOW, True)

tbl_shape = sl.shapes.add_table(5, 9, Inches(0.5), Inches(1.75), Inches(9), Inches(1.85))
t = tbl_shape.table
for j, h in enumerate(["", "0", "1", "2", "3", "4", "5", "6", "7"]):
    c = t.cell(0, j); c.text = h; p = c.text_frame.paragraphs[0]
    p.font.size = Pt(8); p.font.color.rgb = DIM; p.font.name = "Consolas"
    c.fill.solid(); c.fill.fore_color.rgb = BG_INNER

row_data = [
    (["Token","I","went","to","बाज़ार","में","to","buy","groceries"], [DIM]+[WHITE]*8),
    (["Lang","EN","EN","EN","HI","HI","EN","EN","EN"], [DIM]+[CYAN,CYAN,CYAN,ORANGE,ORANGE,CYAN,CYAN,CYAN]),
    (["y_sw","0","0","1","0","1","0","0","—"], [TEAL]+[TEAL,TEAL,RED,TEAL,RED,TEAL,TEAL,DIM]),
    (["y_dur","—","—","S","—","S","—","—","—"], [PURPLE]+[DIM,DIM,PURPLE,DIM,PURPLE,DIM,DIM,DIM]),
]
for ri, (vals, cols) in enumerate(row_data):
    for j, (v, c) in enumerate(zip(vals, cols)):
        cell = t.cell(ri+1, j); cell.text = v
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(8); p.font.color.rgb = c
        p.font.name = "Consolas"; p.font.bold = (j == 0 or v in ["1", "S"])
        cell.fill.solid(); cell.fill.fore_color.rgb = BG_INNER if ri % 2 == 0 else BG_CARD

# Task cards
card(sl, 0.4, 3.95, 4.4, 2.0, BG_CARD)
txt(sl, 0.6, 4.05, 4, 0.25, "Task 1: Switch Prediction (y_sw)", 13, TEAL, True)
for i, line in enumerate(["Binary: will token t+1 switch language?",
                           "0 = no switch,  1 = switch coming",
                           "Last token → -100 (masked in loss)",
                           "~81% no-switch → need weighted loss"]):
    teal_bullet(sl, 0.7, 4.45 + i * 0.32, line, 9)

card(sl, 5.2, 3.95, 4.4, 2.0, BG_CARD)
txt(sl, 5.4, 4.05, 4, 0.25, "Task 2: Duration Prediction (y_dur)", 13, PURPLE, True)
for i, line in enumerate(["3-class: how long will the switch last?",
                           "Small(1-2) · Medium(3-6) · Large(7+)",
                           "Only defined where y_sw = 1",
                           "Non-switch positions → -100 (masked)"]):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(4.52 + i*0.32), Inches(0.1), Inches(0.1))
    s.fill.solid(); s.fill.fore_color.rgb = PURPLE; s.line.fill.background()
    txt(sl, 5.7, 4.45 + i * 0.32, 3.8, 0.25, line, 9, GRAY)

footer(sl, 6)


# =============================================================================
# SLIDE 7: CAUSAL MASK
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)

txt(sl, 0.5, 0.4, 9, 0.6, "Causal Attention Mask", 28, WHITE, True)
txt(sl, 0.5, 0.95, 9, 0.35, "Enforcing prefix-only context — token t can only see tokens 0…t", 12, DIM)

# Bidirectional
card(sl, 0.4, 1.5, 4.4, 2.9, BG_CARD, RED)
txt(sl, 0.6, 1.6, 4, 0.3, "❌  Standard Bidirectional (BERT)", 13, RED, True)
txt(sl, 0.6, 1.95, 4, 0.2, "Every token sees every other token", 9, DIM)
card(sl, 1.0, 2.3, 2.8, 1.5, BG_INNER)
mask1 = "  1   1   1   1   1\n  1   1   1   1   1\n  1   1   1   1   1\n  1   1   1   1   1\n  1   1   1   1   1"
txt(sl, 1.1, 2.35, 2.6, 1.4, mask1, 11, TEAL, font="Consolas", align=PP_ALIGN.CENTER)
txt(sl, 0.6, 3.9, 4, 0.25, "⚠ Token 2 can see token 3 → INFO LEAKAGE", 9, RED, align=PP_ALIGN.CENTER)

# Causal
card(sl, 5.2, 1.5, 4.4, 2.9, BG_CARD, TEAL)
txt(sl, 5.4, 1.6, 4, 0.3, "✅  Causal Mask (Our Approach)", 13, TEAL, True)
txt(sl, 5.4, 1.95, 4, 0.2, "Token t only sees tokens 0…t", 9, DIM)
card(sl, 5.8, 2.3, 2.8, 1.5, BG_INNER)
mask2 = "  1   0   0   0   0\n  1   1   0   0   0\n  1   1   1   0   0\n  1   1   1   1   0\n  1   1   1   1   1"
txt(sl, 5.9, 2.35, 2.6, 1.4, mask2, 11, TEAL, font="Consolas", align=PP_ALIGN.CENTER)
txt(sl, 5.4, 3.9, 4, 0.25, "✓ No future information leakage", 9, TEAL, align=PP_ALIGN.CENTER)

# Implementation
card(sl, 0.4, 4.7, 9.2, 1.5, BG_CARD)
txt(sl, 0.6, 4.8, 4, 0.25, "Implementation (3 lines)", 13, TEAL, True)
card(sl, 0.6, 5.15, 8.8, 0.9, BG_INNER)
code = ("# 1. Lower-triangular mask\n"
        "causal = torch.tril(torch.ones(seq_len, seq_len))\n\n"
        "# 2. Combine with padding mask\n"
        "combined = causal * attention_mask\n\n"
        "# 3. Convert: 0→attend, -10000→block\n"
        "extended = (1.0 - combined) * -10000.0")
txt(sl, 0.7, 5.18, 8.6, 0.85, code, 8, GRAY, font="Consolas")

footer(sl, 7)


# =============================================================================
# SLIDE 8: BASELINE METRICS
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)

txt(sl, 0.5, 0.4, 9, 0.6, "Naive Baseline & Anticipatory F1", 28, WHITE, True)
txt(sl, 0.5, 0.95, 9, 0.35, "Actual results from our pipeline — the floor our model must beat", 12, DIM)

# Left: baseline results
card(sl, 0.4, 1.5, 4.4, 3.3, BG_CARD)
txt(sl, 0.6, 1.6, 4, 0.3, "Naive: 'Always Predict No Switch'", 13, YELLOW, True)
txt(sl, 0.6, 1.95, 4, 0.2, "Predicts class 0 for every single token", 9, DIM)

card(sl, 0.6, 2.3, 4, 2.1, BG_INNER)
metrics = [("Macro F1:", "0.4473", YELLOW), ("F1 (switch=1):", "0.0000", RED),
           ("F1 (no_switch=0):", "0.8947", TEAL), ("", "", WHITE),
           ("Switch rate:", "19.1%", GRAY), ("Class imbalance:", "1:4", GRAY),
           ("Duration: Small", "51.0%", GRAY), ("Duration: Medium", "25.3%", GRAY),
           ("Duration: Large", "23.7%", GRAY)]
for i, (lbl, val, col) in enumerate(metrics):
    if lbl == "": continue
    y = 2.35 + i * 0.23
    txt(sl, 0.7, y, 2.2, 0.2, lbl, 9, GRAY)
    txt(sl, 3.0, y, 1.4, 0.2, val, 9, col, True)

txt(sl, 0.6, 4.5, 4, 0.2, "→ F1(switch)=0 proves baseline is blind to switches!", 9, YELLOW)

# Right: per-pair F1
card(sl, 5.2, 1.5, 4.4, 3.3, BG_CARD)
txt(sl, 5.4, 1.6, 4, 0.3, "Per-Pair F1 & Universality", 13, TEAL, True)

card(sl, 5.4, 2.0, 4, 2.3, BG_INNER)
pairs = [("Hindi-EN", "0.4777", CYAN), ("Arabic-EN", "0.4665", CYAN),
         ("Spanish-EN", "0.4147", ORANGE), ("French-EN", "0.4162", ORANGE),
         ("Korean-EN", "0.4612", PURPLE), ("Chinese-EN", "0.4453", PURPLE)]
for i, (pair, f1, col) in enumerate(pairs):
    y = 2.1 + i * 0.3
    txt(sl, 5.5, y, 2.5, 0.25, pair, 10, col)
    # F1 bar
    bar_w = float(f1) * 3.2
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(y + 0.03), Inches(bar_w), Inches(0.18))
    s.fill.solid(); s.fill.fore_color.rgb = TEAL_DARK; s.line.fill.background(); s.adjustments[0] = 0.3
    txt(sl, 7.3, y, 1.5, 0.25, f1, 9, TEAL, True)

# Sigma
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.95), Inches(3.7), Inches(0.01))
s.fill.solid(); s.fill.fore_color.rgb = TEAL_DARK; s.line.fill.background()

txt(sl, 5.5, 4.0, 2, 0.25, "σ (universality):", 10, WHITE, True)
txt(sl, 7.8, 4.0, 1.2, 0.25, "0.0242", 12, TEAL, True)

# Goal bar
card(sl, 0.4, 5.1, 9.2, 0.6, BG_CARD, TEAL_DARK)
_, tf = txt(sl, 0.6, 5.2, 8.8, 0.4, "", 9)
r = tf.paragraphs[0].add_run(); r.text = "🎯 Goal: "; r.font.color.rgb = TEAL; r.font.bold = True; r.font.size = Pt(10)
r2 = tf.paragraphs[0].add_run()
r2.text = "Trained model must beat 0.4473 macro F1 while maintaining low σ across all 6 language pairs."
r2.font.color.rgb = GRAY; r2.font.size = Pt(9)

footer(sl, 8)


# =============================================================================
# SLIDE 9: CHALLENGES & NEXT STEPS
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); bar(sl)
circle_deco(sl, 8, 5, 3, RGBColor(0x0A, 0x2A, 0x28))

txt(sl, 0.5, 0.4, 9, 0.6, "Challenges & Next Steps", 28, WHITE, True)
txt(sl, 0.5, 0.95, 9, 0.35, "Sub-word tokenization issues and the Phase 2 roadmap", 12, DIM)

# Challenges
card(sl, 0.4, 1.4, 4.4, 4.6, BG_CARD, RED)
txt(sl, 0.6, 1.5, 4, 0.3, "⚠  Current Challenges", 14, RED, True)

challenges = [
    ("Subword Misalignment", '"supermercado" → ["▁super","merc","ado"]', "→ Predict only on first subword, mask rest"),
    ("Same-Script LID", '"hospital" exists in both EN and ES', "→ Using langid classifier for Latin↔Latin"),
    ("Chinese-EN Shorter Texts", "27.5 avg words vs 45-57 for other pairs", "→ Monitor impact on prediction quality"),
    ("Class Imbalance (1:4)", "81% of tokens are 'no-switch'", "→ Weighted CrossEntropy or focal loss"),
]
for i, (title, desc, fix) in enumerate(challenges):
    y = 1.95 + i * 0.98
    card(sl, 0.6, y, 4, 0.85, BG_INNER)
    txt(sl, 0.75, y + 0.05, 3.7, 0.2, title, 10, WHITE, True)
    txt(sl, 0.75, y + 0.28, 3.7, 0.2, desc, 8, DIM)
    txt(sl, 0.75, y + 0.5, 3.7, 0.2, fix, 8, TEAL)

# Next steps
card(sl, 5.2, 1.4, 4.4, 4.6, BG_CARD, TEAL)
txt(sl, 5.4, 1.5, 4, 0.3, "🚀  Next Steps (Weeks 5-7)", 14, TEAL, True)

steps = [
    ("Week 5", "Build Dual-Head Model", ["XLM-RoBERTa encoder + causal mask", "Head 1: switch prediction (binary)", "Head 2: duration prediction (3-class)"]),
    ("Week 6", "Training Pipeline", ["Combined loss: λ₁·L_sw + λ₂·L_dur", "Mixed precision (fp16) training", "Early stopping on validation F1"]),
    ("Week 7", "Architecture Comparison", ["Train both XLM-RoBERTa and mBERT", "Compare per-pair F1 scores", "Compute universality σ for each"]),
]
for i, (week, task, items) in enumerate(steps):
    y = 1.95 + i * 1.35
    card(sl, 5.4, y, 4, 1.2, BG_INNER)
    # Week badge
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(y + 0.08), Inches(0.7), Inches(0.22))
    s.fill.solid(); s.fill.fore_color.rgb = TEAL_DARK; s.line.fill.background(); s.adjustments[0] = 0.3
    txt(sl, 5.5, y + 0.07, 0.7, 0.22, week, 8, TEAL, True, PP_ALIGN.CENTER)
    txt(sl, 6.3, y + 0.05, 2.8, 0.25, task, 10, WHITE, True)
    for j, item in enumerate(items):
        txt(sl, 5.6, y + 0.38 + j * 0.23, 3.6, 0.2, f"•  {item}", 8, GRAY)

footer(sl, 9)


# =============================================================================
# SAVE
# =============================================================================
out = "update1_presentation.pptx"
prs.save(out)
print(f"\n✅ Saved: {os.path.abspath(out)}")
print("   Open in PowerPoint or upload to Google Slides!")